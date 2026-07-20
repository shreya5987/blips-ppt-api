import json
import argparse
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Color palette
# -----------------------------

DEFAULT_TEMPLATE_PATH = Path("templates/AI Agent testing - WET TEMPLATE.pptx")

COLORS = {
    "dark_blue": RGBColor(31, 78, 121),
    "light_blue": RGBColor(231, 245, 255),
    "gray": RGBColor(89, 89, 89),
    "light_gray": RGBColor(242, 242, 242),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "green": RGBColor(112, 173, 71),
    "yellow": RGBColor(255, 220, 52),
    "red": RGBColor(192, 0, 0),
    "purple": RGBColor(67, 0, 122),
    "light_yellow": RGBColor(255, 245, 209),
    "orange": RGBColor(255, 96, 18)
}



# Formatting helpers
# -----------------------------


def set_text_box_text(
    shape,
    text,
    font_size=16,
    bold=True,
    color=COLORS["black"],
    align=PP_ALIGN.LEFT,
):
    """
    Safely set text in a textbox/placeholder/shape.
    """
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    run.text = str(text)



    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color




def add_footer(slide, footer_text="Generated from JSON analysis"):
    """
    Add small footer.
    """
    shape = slide.shapes.add_textbox(
        Inches(7.8),
        Inches(5.35),
        Inches(1.75),
        Inches(0.25),
    )

    set_text_box_text(
        shape,
        footer_text,
        font_size=8,
        bold=False,
        color=COLORS["black"],
    )

    return shape



def add_section_box(
    slide,
    title,
    body,
    left,
    top,
    width,
    height,
    title_color=COLORS["purple"],
    body_font_size=13,
):
    """
    Add a rounded rectangle section with title and body.
    """

    # Outer box
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )

    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["white"]
    box.line.color.rgb = COLORS["light_blue"]

    # Section title
    title_shape = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.1),
        width - Inches(0.3),
        Inches(0.35),
    )

    set_text_box_text(
        title_shape,
        title,
        font_size=14,
        bold=True,
        color=title_color,
    )

    # Section body
    body_shape = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.6),
    )

    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True

    lines = str(body).split("\n")

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)

        run = p.add_run()
        run.text = line
        run.font.size = Pt(body_font_size)
        run.font.color.rgb = COLORS["black"]

    return body_shape



def add_title(slide, title_text):
    """
    Add a consistent slide title.
    """
    left = Inches(0.4)
    top = Inches(0.25)
    width = Inches(9.0)
    height = Inches(0.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    set_text_box_text(
        shape,
        title_text,
        font_size=24,
        bold=True,
        color=COLORS["purple"],
    )

    # Set title font

    for paragraph in shape.text_frame.paragraphs:

        for run in paragraph.runs:

            run.font.name = "Arial Black"
    
    return shape




def add_actions_table(slide, actions, left, top, width, height):
    """
    Add action table.
    """

    if not isinstance(actions, list):
        actions = []

    rows = max(len(actions), 1) + 1
    cols = 2

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        height,
    )

    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(2.0)         #changed from 2.2 to 1.8
    table.columns[1].width = Inches(7.0)        #changed from 7.7 to 6.0

    headers = ["Owner", "Action"]

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["purple"]

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(15)
                run.font.color.rgb = COLORS["white"]

    # Body rows
    if len(actions) == 0:
        actions = [
            {
                "owner": "",
                "action": "No actions provided",
                
            }
        ]

    for row_idx, action_item in enumerate(actions, start=1):
        values = [
            action_item.get("owner", ""),
            action_item.get("action", ""),
    
        ]

        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light_gray"]

            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light_blue"]


            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if col_idx != 2 else PP_ALIGN.CENTER
                paragraph.space_after = Pt(2)

                for run in paragraph.runs:
                    run.font.size = Pt(12)     #changed from 9 to 12
                    run.font.color.rgb = COLORS["black"]

    return table_shape


# -----------------------------
# Template/layout helper
# -----------------------------

def get_blank_layout(prs):
    """
    Try to find a blank layout from the template.
    Falls back to last available layout.
    """
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout

    return prs.slide_layouts[-1]


# -----------------------------
# Slide builders
# -----------------------------

def build_title_slide(prs, data):
    """
    Build executive summary slide.
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))

    add_title(slide, data.get("title"))

    subtitle = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.85),
        Inches(8.0),
        Inches(0.4),
    )

    set_text_box_text(
        subtitle,
        data.get("subtitle"),
        font_size=16,
        bold=True,
        color=COLORS["gray"],
    )

    add_section_box(
        slide,
        "Executive Summary",
        data.get("summary"),
        Inches(0.55),
        Inches(1.45),
        Inches(8.7),      #changed from 11.9 to 8.5
        Inches(2.3),       #changed from 3.5 to 2.5
        body_font_size=13,       #changed from 13 to 14
    )

    

    add_section_box(
        slide,
        "Primary Failure Interpretation",
        data.get("primaryFailure"),
        Inches(0.55),
        Inches(4.0),  #changed from 5.15 to 
        Inches(8.7),               #changed from 8.5 to 5.5
        Inches(1.3),
        body_font_size=13,
    )

    add_footer(slide)

    return slide


def build_recommendation_slide(prs, data):
    """
    Build recommendation slide.
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))

    add_title(slide, "Recommended Investigation Path")

    add_section_box(
        slide,
        "Recommendation",
        data.get("recommendation", ""),
        Inches(0.55),
        Inches(1.05),
        Inches(8.7),      #changed from 11.9 to 8.5
        Inches(2.0),         #changed from 4.0 to 2.0
        body_font_size=14,       #changed from 13 to 14
    )

    
    # new section for key focus areas to be iterative
    focus_data = data.get("keyFocusAreas", "")

    if isinstance(focus_data, list):
        # Convert list → numbered string
        focus_areas = "\n".join(
            f"{i+1}. {item}" for i, item in enumerate(focus_data)
        )
    else:
        # Already a string
        focus_areas = str(focus_data)


    add_section_box(
        slide,
        "Key Focus Areas",
        focus_areas,
        Inches(0.55),
        Inches(3.3),       #changed from 5.3 to 3.5
        Inches(8.7),      #changed from 11.9 to 8.5
        Inches(2),    #changed from 1.35 to 2
        body_font_size=12,
    )
    add_footer(slide)

    return slide


def build_actions_slide(prs, data):
    """
    Build actions table slide.
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))

    add_title(slide, "Action Plan")

    actions = data.get("actions", [])

    add_actions_table(
        slide,
        actions,
        Inches(0.45),
        Inches(1.1),
        Inches(3),      #changed from 12.3 to 10.3
        Inches(4.2),      #changed from 5.5 to 4.2
    )

    add_footer(slide)

    return slide



def build_bibliography_slide(prs, data):
    """
    Build bibliography slide.
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))

    add_title(slide, "References")

    
    references_data = data.get("references", [])

    ref_box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.6),
        Inches(9.2),
        Inches(4.7),
    )

    # Get text frame
    text_frame = ref_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True


    for i, ref in enumerate(references_data, start=1):
        p = text_frame.add_paragraph()

        # Number (plain text)
        run_num = p.add_run()
        run_num.text = f"{i}. "

        # Hyperlinked title
        run_link = p.add_run()
        run_link.text = ref.get("text", "Untitled")
        run_link.font.size = Pt(13)

        if ref.get("url"):
            run_link.hyperlink.address = ref["url"]

        desc = ref.get("description", "")

        if desc:
            p_desc = text_frame.add_paragraph()

            run_desc = p_desc.add_run()

            run_desc.text = desc
            p_desc.left_margin = Inches(0.3)
            run_desc.font.size = Pt(13)

        text_frame.add_paragraph()  # Add a blank line after each reference


    add_footer(slide)

    return slide






def move_appended_slides_after(prs, generated_start_index, after_slide_number=3):
    """
    Moves all slides added after generated_start_index to appear after a given slide number.

    after_slide_number is 1-based.
    Example: after_slide_number=3 means insert after slide 3.
    """

    if len(prs.slides) < after_slide_number:
        raise ValueError(f"Template must have at least {after_slide_number} slides.")

    sldIdLst = prs.slides._sldIdLst

    # These are the slides that were generated and appended to the end
    generated_slide_ids = list(sldIdLst)[generated_start_index:]

    # Remove generated slides from the end
    for sldId in generated_slide_ids:
        sldIdLst.remove(sldId)

    # Insert after slide 3
    insert_index = after_slide_number

    for offset, sldId in enumerate(generated_slide_ids):
        sldIdLst.insert(insert_index + offset, sldId)








# -----------------------------
# Main Jupyter function
# -----------------------------

def json_to_powerpoint_notebook(
    json_path,
    template_path,
    output_path="AI Agent testing - WET TEMPLATE.pptx",
    keep_existing_template_slides=True,
):
    """
    Convert an RCA-style JSON file into a PowerPoint deck using a PPTX template.

    Parameters
    ----------
    json_path : str or Path
        Path to input JSON file.

    template_path : str or Path
        Path to PowerPoint template .pptx file.

    output_path : str or Path, optional
        Output PowerPoint file path.

    keep_existing_template_slides : bool, optional
        If False, existing slides in the template are removed.
        If True, generated slides are appended after existing template slides.

    Returns
    -------
    output_path : Path
        Path to the saved PowerPoint deck.
    """

    json_path = Path(json_path)
    template_path = Path(template_path)
    output_path = Path(output_path)

    # Validate files
    if not json_path.exists():
        raise FileNotFoundError(f"JSON file not found: {json_path.resolve()}")

    if not template_path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path.resolve()}")

    # Load JSON
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    # Validate required keys
    required_keys = ["summary", "recommendation", "actions"]
    missing = [k for k in required_keys if k not in data]

    if missing:
        raise ValueError(f"JSON missing required key(s): {missing}")

    # Load PowerPoint template
    prs = Presentation(str(template_path))

    # Optional: remove existing template slides
    if not keep_existing_template_slides:
        while len(prs.slides) > 0:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]


    generated_start_index = len(prs.slides)

    # Build slides
    build_title_slide(prs, data)
    build_recommendation_slide(prs, data)
    build_actions_slide(prs, data)
    build_bibliography_slide(prs, data)

        
    #move_appended_slides_after(
     #   prs,
      #  generated_start_index,
       # after_slide_number=3,
        #)

    # Save deck
    prs.save(str(output_path))

    print(f"PowerPoint deck created: {output_path.resolve()}")

    return output_path


def generate_ppt(
data: dict,
output_path: str,
template_path: str | Path = "templates/AI Agent testing - WET TEMPLATE.pptx",
) -> str:
    """
    FastAPI-friendly wrapper.

    Takes JSON data as a Python dictionary, writes it to a temporary JSON file,
    then calls the existing json_to_powerpoint_notebook function.
    """

    with tempfile.NamedTemporaryFile(
        mode="w",
        suffix=".json",
        delete=False,
        encoding="utf-8",
    ) as temp_json:
        json.dump(data, temp_json, indent=2)
        temp_json_path = Path(temp_json.name)

    try:
        generated_path = json_to_powerpoint_notebook(
            json_path=temp_json_path,
            template_path=template_path,
            output_path=output_path,
            keep_existing_template_slides=True,
        )

        return str(generated_path)

    finally:
        if temp_json_path.exists():
            temp_json_path.unlink()
